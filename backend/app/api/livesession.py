from fastapi import APIRouter, Depends, HTTPException, UploadFile, File, Form, status, Body
from typing import Dict, Any, List, Optional
from pydantic import BaseModel, ConfigDict
import logging
import io
import fitz  # PyMuPDF
from pptx import Presentation

from app.services.live_session_service import LiveSessionService
from app.api.auth import get_current_user
from app.db import get_db
from datetime import datetime
from bson import ObjectId

logger = logging.getLogger(__name__)
router = APIRouter()
ls_service = LiveSessionService()

class GenerateRequest(BaseModel):
    content: str
    level: int
    subject_area: str
    topic: str
    mcq_count: int
    short_count: int
    coding_count: int

class GradeRequest(BaseModel):
    model_config = ConfigDict(protected_namespaces=())
    question: str
    model_answer: str
    keywords: List[str]
    student_answer: str
    max_score: int
    level: int

def extract_text_from_file(file_content: bytes, filename: str) -> str:
    text = ""
    try:
        if filename.lower().endswith(".pdf"):
            doc = fitz.open(stream=file_content, filetype="pdf")
            for page in doc:
                text += page.get_text("text") + "\n"
        elif filename.lower().endswith((".ppt", ".pptx")):
            prs = Presentation(io.BytesIO(file_content))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        else:
            text = file_content.decode("utf-8", errors="ignore")
    except Exception as e:
        logger.error(f"Text extraction failed: {e}")
        raise ValueError("Unsupported or corrupt file format")
        
    return text[:8000] # Cap at requested limits

@router.post("/upload")
async def upload_document(
    file: UploadFile = File(...),
    subject_area: str = Form("General Education"),
    topic_hint: str = Form(""),
    current_user=Depends(get_current_user)
):
    try:
        content_bytes = await file.read()
        extracted_text = extract_text_from_file(content_bytes, file.filename)
        
        if len(extracted_text.split()) < 100:
            return {
                "error": "INSUFFICIENT_CONTENT", 
                "message": "The uploaded file did not contain enough text to generate questions. Please upload a more detailed document or add a topic hint."
            }

        summary_data = await ls_service.summarize(extracted_text, subject_area, topic_hint)
        
        return {
            "success": True,
            "raw_content": extracted_text,
            "summary": summary_data
        }
    except ValueError as ve:
        raise HTTPException(status_code=400, detail=str(ve))
    except Exception as e:
        logger.error(f"Upload failed: {e}")
        raise HTTPException(status_code=500, detail="Failed to process document")

@router.post("/generate")
async def generate_questions(
    req: GenerateRequest,
    current_user=Depends(get_current_user)
):
    try:
        questions = await ls_service.generate_questions(
            content=req.content,
            level=req.level,
            subject_area=req.subject_area,
            topic=req.topic,
            mcq_count=req.mcq_count,
            short_count=req.short_count,
            coding_count=req.coding_count
        )
        return {"success": True, "questions": questions}
    except Exception as e:
        logger.error(f"Generate failed: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/grade")
async def grade_answer(
    req: GradeRequest,
    current_user=Depends(get_current_user)
):
    try:
        result = await ls_service.grade_short_answer(
            question=req.question,
            model_answer=req.model_answer,
            keywords=req.keywords,
            student_answer=req.student_answer,
            max_score=req.max_score,
            level=req.level
        )
        return {"success": True, "feedback": result}
    except Exception as e:
        logger.error(f"Grade failed: {e}")
        raise HTTPException(status_code=500, detail=str(e))


# ──────────────────────────────────────────────
# Live Classroom Session Management
# ──────────────────────────────────────────────

class SessionStartRequest(BaseModel):
    batch_id: str
    topic: Optional[str] = ""

@router.post("/sessions/start")
async def start_live_session(
    req: SessionStartRequest,
    current_user=Depends(get_current_user)
):
    """Create a new live classroom session record for a batch."""
    try:
        db = await get_db()

        user_id = str(current_user.id) if hasattr(current_user, 'id') else str(current_user.get("_id", ""))

        session_doc = {
            "batch_id": req.batch_id,
            "teacher_id": user_id,
            "topic": req.topic,
            "started_at": datetime.utcnow(),
            "status": "active",
            "students_joined": [],
        }
        result = await db.live_sessions.insert_one(session_doc)
        logger.info(f"[LIVESESSION] Session started for batch {req.batch_id}: {result.inserted_id}")

        return {
            "success": True,
            "session_id": str(result.inserted_id),
            "batch_id": req.batch_id,
            "status": "active"
        }
    except Exception as e:
        logger.error(f"[LIVESESSION] Session start failed: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to start session: {str(e)}")


@router.post("/sessions/{session_id}/end")
async def end_live_session(
    session_id: str,
    current_user=Depends(get_current_user)
):
    """Mark a live session as ended."""
    try:
        db = await get_db()
        await db.live_sessions.update_one(
            {"_id": ObjectId(session_id)},
            {"$set": {"status": "ended", "ended_at": datetime.utcnow()}}
        )
        return {"success": True, "session_id": session_id}
    except Exception as e:
        logger.error(f"[LIVESESSION] Session end failed: {e}")
        raise HTTPException(status_code=500, detail=str(e))

