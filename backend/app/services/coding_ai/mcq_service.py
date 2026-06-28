from datetime import timezone
from google import genai
from google.genai import types
import json
import os
import re
import asyncio
from typing import Dict, List, Any, Optional, Tuple
from datetime import datetime
from dotenv import load_dotenv

load_dotenv()


class MCQServiceMixin:
    async def generate_mcq_questions(
            self, 
            topic: str, 
            difficulty: str, 
            count: int = 10,
            store_in_db: bool = True
        ) -> List[Dict[str, Any]]:
            """Generate MCQ questions using Gemini AI with automatic model fallback."""
            if not self.available:
                print(f"[GEMINI] Client unavailable — serving fallback questions for '{topic}'")
                return self._get_fallback_mcq_questions(topic, difficulty, count)

            # Get existing questions to build uniqueness hint
            try:
                existing_questions = await self._get_existing_questions(topic, difficulty)
                existing_texts = [q.get("question", "") for q in existing_questions[:5]]
            except Exception:
                existing_texts = []

            uniqueness_hint = ""
            if existing_texts:
                uniqueness_hint = "\nAvoid questions similar to:\n" + "\n".join(f"- {q}" for q in existing_texts)

            prompt = f"""Generate {count} UNIQUE multiple choice questions about {topic} at {difficulty} difficulty.

    Rules:
    - Exactly 4 options per question (A, B, C, D)
    - Educational, clear, and unambiguous
    - Cover varied subtopics and question formats
    - Distribute correct answers across A, B, C, D (don't always use A)
    - explanation must justify why the chosen option is correct
    {uniqueness_hint}

    Return ONLY a valid JSON array:
    [
      {{
        "question": "What does the 'def' keyword do in Python?",
        "options": ["Defines a variable", "Defines a function", "Defines a class", "Defines a module"],
        "answer": "B",
        "explanation": "Defines a function is correct because 'def' is the Python keyword used to declare functions."
      }}
    ]

    Return ONLY the JSON array, no markdown, no extra text."""

            # Build model priority: active model first, then rest of fallback list
            models_to_try = [self.active_model] + [
                m for m in self.fallback_models if m != self.active_model
            ]

            response = None
            last_error = None

            for model in models_to_try:
                try:
                    print(f"[GEMINI] Generating {count} questions for '{topic}' using {model}")
                    response = await self.client.aio.models.generate_content(
                        model=model,
                        contents=prompt,
                        config=types.GenerateContentConfig(
                            temperature=0.7,
                            max_output_tokens=4096
                        )
                    )
                    if model != self.active_model:
                        print(f"[GEMINI] Switched active model: {self.active_model} → {model}")
                        self.active_model = model
                    break  # success
                except Exception as err:
                    err_str = str(err).lower()
                    is_capacity = any(kw in err_str for kw in [
                        "429", "resource_exhausted", "quota", "rate limit",
                        "too many requests", "overloaded", "high demand",
                        "service unavailable", "503"
                    ])
                    if is_capacity:
                        print(f"[GEMINI] {model} overloaded — trying next model...")
                        last_error = err
                        continue
                    # Unknown error on this model — try next model anyway
                    print(f"[GEMINI] {model} error: {err} — trying next model...")
                    last_error = err
                    continue

            if response is None:
                print(f"[GEMINI] All {len(models_to_try)} models failed — serving fallback questions. Last error: {last_error}")
                return self._get_fallback_mcq_questions(topic, difficulty, count)

            # Parse response
            try:
                content = self._clean_json_response(response.text.strip())
                start = content.find('[')
                end = content.rfind(']')
                if start != -1 and end != -1:
                    content = content[start:end + 1]

                questions = json.loads(content)
                formatted = []

                for i, q in enumerate(questions):
                    if not all(k in q for k in ["question", "options"]):
                        continue

                    raw_answer = q.get("answer", q.get("correct_answer", "A"))
                    correct_idx = 0
                    if isinstance(raw_answer, int):
                        correct_idx = max(0, min(3, raw_answer))
                    elif isinstance(raw_answer, str):
                        clean = raw_answer.upper().strip()
                        if clean and clean[0] in "ABCD":
                            correct_idx = ord(clean[0]) - ord("A")

                    options = q.get("options", [])
                    correct_text = options[correct_idx] if 0 <= correct_idx < len(options) else ""

                    formatted_q = {
                        "id": f"q{i+1}",
                        "question": q["question"],
                        "options": options,
                        "answer": correct_text,
                        "correct_answer": correct_idx,
                        "explanation": q.get("explanation", ""),
                        "difficulty": difficulty,
                        "topic": topic,
                        "generated_by": "gemini",
                        "type": "mcq"
                    }

                    if self._validate_mcq_question(formatted_q):
                        formatted.append(formatted_q)

                if not formatted:
                    print(f"[GEMINI] No valid questions parsed from response — serving fallback")
                    return self._get_fallback_mcq_questions(topic, difficulty, count)

                print(f"[GEMINI] Generated {len(formatted[:count])}/{count} questions for '{topic}' ✓")

                if store_in_db:
                    try:
                        await self._store_ai_questions_in_db(formatted[:count], topic, difficulty)
                    except Exception:
                        pass  # DB storage failure is non-critical

                return formatted[:count]

            except (json.JSONDecodeError, Exception) as parse_err:
                print(f"[GEMINI] JSON parse failed: {parse_err} — serving fallback questions")
                return self._get_fallback_mcq_questions(topic, difficulty, count)

    def _get_fallback_mcq_questions(self, topic: str, difficulty: str, count: int) -> List[Dict[str, Any]]:
            """Get fallback MCQ questions when AI is not available - with enhanced variety"""
            import random
            from datetime import datetime
            
            print(f" [GEMINI] Using fallback questions for {topic} ({difficulty})")
            
            # Expanded question pool with much more variety
            fallback_questions = []
            
            # Comprehensive topic-specific question pools
            topic_questions = {
                "python": [
                    {
                        "question": "What is the correct syntax to create a list in Python?",
                        "options": ["list = []", "list = {}", "list = ()", "list = []"],
                        "correct": 0,
                        "explanation": "'list = []' is correct because square brackets [] are the proper syntax for creating lists in Python."
                    },
                    {
                        "question": "Which keyword is used to define a function in Python?",
                        "options": ["function", "def", "define", "func"],
                        "correct": 1,
                        "explanation": "'def' is correct because it is the specific keyword used to define functions in Python."
                    },
                    {
                        "question": "What does the 'print()' function do in Python?",
                        "options": ["Displays output", "Reads input", "Calculates values", "Stores data"],
                        "correct": 0,
                        "explanation": "'Displays output' is correct because the print() function outputs text to the console."
                    },
                    {
                        "question": "Which data type is used to store a sequence of characters in Python?",
                        "options": ["string", "int", "float", "boolean"],
                        "correct": 0,
                        "explanation": "'string' is correct because it is the data type used to store sequences of characters in Python."
                    },
                    {
                        "question": "What is the result of 5 // 2 in Python?",
                        "options": ["2.5", "2", "3", "2.0"],
                        "correct": 1,
                        "explanation": "The // operator performs floor division, returning 2"
                    },
                    {
                        "question": "What does the len() function return?",
                        "options": ["The last element", "The length of a sequence", "The sum of elements", "The average"],
                        "correct": 1,
                        "explanation": "len() returns the number of items in a sequence or collection"
                    },
                    {
                        "question": "Which method is used to add an element to the end of a list?",
                        "options": ["add()", "append()", "insert()", "extend()"],
                        "correct": 1,
                        "explanation": "append() adds a single element to the end of a list"
                    },
                    {
                        "question": "What is the result of 5 % 2 in Python?",
                        "options": ["2.5", "1", "2", "0"],
                        "correct": 1,
                        "explanation": "The % operator returns the remainder of division, so 5 % 2 = 1"
                    }
                ],
                "javascript": [
                    {
                        "question": "What is the correct way to declare a variable in JavaScript?",
                        "options": ["var x = 5", "variable x = 5", "v x = 5", "declare x = 5"],
                        "correct": 0,
                        "explanation": "Variables in JavaScript are declared using 'var', 'let', or 'const'"
                    },
                    {
                        "question": "Which operator is used for strict equality in JavaScript?",
                        "options": ["==", "===", "=", "!="],
                        "correct": 1,
                        "explanation": "The === operator checks for strict equality (value and type)"
                    },
                    {
                        "question": "What is the result of typeof null in JavaScript?",
                        "options": ["null", "undefined", "object", "string"],
                        "correct": 2,
                        "explanation": "typeof null returns 'object' due to a historical bug in JavaScript"
                    },
                    {
                        "question": "Which method is used to add an element to the end of an array?",
                        "options": ["push()", "add()", "append()", "insert()"],
                        "correct": 0,
                        "explanation": "push() adds one or more elements to the end of an array"
                    }
                ],
                "array": [
                    {
                        "question": "What is the time complexity of accessing an element by index in an array?",
                        "options": ["O(1)", "O(n)", "O(log n)", "O(n)"],
                        "correct": 0,
                        "explanation": "Array access by index is O(1) because it uses direct memory addressing"
                    },
                    {
                        "question": "What is the space complexity of an array with n elements?",
                        "options": ["O(1)", "O(n)", "O(log n)", "O(n)"],
                        "correct": 1,
                        "explanation": "An array with n elements requires O(n) space to store all elements"
                    },
                    {
                        "question": "What is the time complexity of linear search in an unsorted array?",
                        "options": ["O(1)", "O(log n)", "O(n)", "O(n)"],
                        "correct": 2,
                        "explanation": "Linear search checks each element sequentially, taking O(n) time in worst case"
                    },
                    {
                        "question": "What is the time complexity of binary search in a sorted array?",
                        "options": ["O(1)", "O(log n)", "O(n)", "O(n)"],
                        "correct": 1,
                        "explanation": "Binary search eliminates half the search space each iteration, taking O(log n) time"
                    }
                ]
            }
            
            # Get questions for the topic or use a default set
            questions = topic_questions.get(topic.lower(), topic_questions["python"])
            
            # If we need more questions than available, expand the pool
            if count > len(questions):
                # Duplicate and modify questions to create more variety
                expanded_questions = questions.copy()
                for i in range(count - len(questions)):
                    base_question = questions[i % len(questions)]
                    # Create variation by modifying the question slightly
                    modified_question = base_question.copy()
                    modified_question["question"] = f"{base_question['question']} (Variant {i+1})"
                    expanded_questions.append(modified_question)
                questions = expanded_questions
            
            # Generate questions with better variety using timestamp-based seeding
            used_indices = set()
            for i in range(count):
                # Use timestamp and topic for better randomization
                random.seed(hash(f"{topic}_{difficulty}_{datetime.now(timezone.utc).timestamp()}_{i}_{len(used_indices)}"))
                
                # Ensure we don't repeat questions
                available_indices = [idx for idx in range(len(questions)) if idx not in used_indices]
                if not available_indices:
                    # If all questions used, reset and continue
                    used_indices.clear()
                    available_indices = list(range(len(questions)))
                
                selected_idx = random.choice(available_indices)
                used_indices.add(selected_idx)
                selected_question = questions[selected_idx]
                
                fallback_questions.append({
                    "id": f"q{i+1}",
                    "question": selected_question["question"],
                    "options": selected_question["options"],
                    "answer": selected_question["options"][selected_question["correct"]],
                    "correct_answer": selected_question["correct"],
                    "explanation": selected_question["explanation"],
                    "difficulty": difficulty,
                    "topic": topic,
                    "type": "mcq"
                })
            
            return fallback_questions

    async def _get_existing_questions(self, topic: str, difficulty: str, limit: int = 20) -> List[Dict[str, Any]]:
            """Get existing questions for a topic and difficulty to avoid duplicates"""
            try:
                from ..db.session import get_db
                
                db = await get_db()
                
                # Query existing questions from multiple collections
                existing_questions = []
                
                # Check ai_questions collection
                ai_questions = await db.ai_questions.find({
                    "topic": {"$regex": topic, "$options": "i"},
                    "difficulty": difficulty,
                    "status": "active"
                }).limit(limit).to_list(length=None)
                existing_questions.extend(ai_questions)
                
                # Check assessments collection for questions
                assessments = await db.assessments.find({
                    "subject": {"$regex": topic, "$options": "i"},
                    "difficulty": difficulty,
                    "questions": {"$exists": True, "$ne": []}
                }).limit(limit).to_list(length=None)
                
                for assessment in assessments:
                    for question in assessment.get("questions", []):
                        if question.get("question"):
                            existing_questions.append(question)
                
                # Check teacher_assessments collection
                teacher_assessments = await db.teacher_assessments.find({
                    "topic": {"$regex": topic, "$options": "i"},
                    "difficulty": difficulty,
                    "questions": {"$exists": True, "$ne": []}
                }).limit(limit).to_list(length=None)
                
                for assessment in teacher_assessments:
                    for question in assessment.get("questions", []):
                        if question.get("question"):
                            existing_questions.append(question)
                
                return existing_questions[:limit]
                
            except Exception as e:
                print(f" [GEMINI] Error getting existing questions: {str(e)}")
                return []

    async def _store_ai_questions_in_db(self, questions: List[Dict[str, Any]], topic: str, difficulty: str):
            """Store AI-generated questions in the database"""
            try:
                from ..db.session import get_db
                from datetime import datetime
                
                db = await get_db()
                
                # Prepare questions for database storage
                questions_to_store = []
                for question in questions:
                    question_doc = {
                        "question": question["question"],
                        "options": question["options"],
                        "answer": question["answer"],
                        "explanation": question["explanation"],
                        "topic": topic,
                        "difficulty": difficulty,
                        "generated_by": "gemini",
                        "metadata": {
                            "ai_model": "gemini-3-flash-preview",
                            "generation_timestamp": datetime.now(timezone.utc).isoformat(),
                            "original_topic": topic,
                            "original_difficulty": difficulty
                        },
                        "created_at": datetime.now(timezone.utc),
                        "status": "active",
                        "usage_count": 0,
                        "quality_score": None
                    }
                    questions_to_store.append(question_doc)
                
                # Insert questions into database
                if questions_to_store:
                    result = await db.ai_questions.insert_many(questions_to_store)
                    print(f" [GEMINI] Stored {len(result.inserted_ids)} AI questions in database")
                    
            except Exception as e:
                print(f" [GEMINI] Error storing AI questions in database: {str(e)}")

    def _validate_mcq_question(self, question: Dict[str, Any]) -> bool:
            """Validate MCQ question structure"""
            required_fields = ["question", "options"]
            if not all(field in question for field in required_fields):
                return False
            
            # Check if we have either "answer" or "correct_answer"
            if "answer" not in question and "correct_answer" not in question:
                return False
            
            if not isinstance(question["options"], list) or len(question["options"]) != 4:
                return False
            
            # If we have correct_answer, validate it's an integer index
            if "correct_answer" in question:
                if not isinstance(question["correct_answer"], int) or question["correct_answer"] < 0 or question["correct_answer"] > 3:
                    return False
            
            return True

    def _generate_fallback_mcq_questions(self, topic: str, difficulty: str, count: int) -> List[Dict[str, Any]]:
            """Generate fallback MCQ questions when AI is not available"""
            print(f" [GEMINI_CODING] Using fallback MCQ generation for {topic}")
            
            fallback_questions = []
            for i in range(count):
                fallback_questions.append({
                    "question": f"What is the main concept of {topic}? (Question {i+1})",
                    "options": [
                        f"Option A for {topic}",
                        f"Option B for {topic}",
                        f"Option C for {topic}",
                        f"Option D for {topic}"
                    ],
                    "correct_answer": i % 4,
                    "explanation": f"This is a fallback question about {topic}",
                    "difficulty": difficulty,
                    "topic": topic,
                    "generated_by": "fallback",
                    "type": "mcq"
                })
            
            return fallback_questions

    async def parse_course_handout(self, file_content: bytes, mime_type: str, subject: str) -> List[Dict[str, Any]]:
            """Parse course handout (PDF/Image/Text) to extract session topics"""
            try:
                print(f" [GEMINI_CODING] Parsing handout for {subject} (Type: {mime_type})")
                
                if not self.available:
                    return [
                        {"topic": "Introduction to " + subject, "description": "Basic concepts"},
                        {"topic": subject + " Fundamentals", "description": "Core principles"},
                        {"topic": "Advanced " + subject, "description": "Complex topics"}
                    ]
                
                prompt_text = f"""
                Analyze the attached Course Handout/Syllabus file for the subject "{subject}" and extract a definitive list of teaching sessions.
                The file likely contains units, modules, or a day-by-day plan.
                
                CRITICAL INSTRUCTION:
                - **GROUP TOPICS**: The syllabus lists many small topics. You MUST group 3-4 adjacent, related topics into a SINGLE session (approx 1 hour).
                - **CONSOLIDATE**: Do NOT create a separate session for every single line item. Example: Instead of 3 sessions for "Intro", "Definition", "Scope", create 1 session "Introduction: Definition & Scope".
                - **COVERAGE**: Ensure the entire syllabus is covered, but in fewer, chunkier sessions.
                
                Return ONLY a valid JSON array of objects with this structure:
                [
                    {{
                        "topic": "Session Title (Consolidated)",
                        "description": "Comma-separated list of sub-topics covered in this session",
                        "unit": "Unit 1" (optional)
                    }}
                ]
                
                Detailed Instructions:
                - Use the visual layout and text to infer sections.
                - Ignore administrative details (policies, grading).
                """
                
                contents = [prompt_text]
                
                # Add file content as a part based on mime type
                if mime_type.startswith("image/") or mime_type == "application/pdf":
                    contents.append({
                        "mime_type": mime_type,
                        "data": file_content
                    })
                else:
                    # Treat as text
                    try:
                        text_data = file_content.decode('utf-8', errors='ignore')
                        contents.append(f"\n\nHandout Content:\n{text_data[:20000]}")
                    except Exception as e:
                        print(f"Error decoding text file: {e}")
                
                import asyncio
                try:
                    # Use new SDK client.aio.models.generate_content for async
                    response = await self.client.aio.models.generate_content(
                        model=self.model_name,
                        contents=contents,
                        config=types.GenerateContentConfig(
                            temperature=0.7,
                            max_output_tokens=4096
                        )
                    )
                except Exception as e:
                    print(f"[ERROR] [GEMINI_CODING] Handout processing failed: {e}")
                    return []
                    
                if not response.text:
                    return []
                    
                json_text = self._clean_json_response(response.text.strip())
                
                try:
                    data = json.loads(json_text)
                    if isinstance(data, list):
                        return data
                    if isinstance(data, dict) and "sessions" in data:
                        return data["sessions"]
                    return []
                except json.JSONDecodeError:
                    print(f"[ERROR] [GEMINI_CODING] JSON parse failed for Handout Parsing")
                    return []

            except Exception as e:
                print(f" [GEMINI_CODING] Error parsing handout: {e}")
                return []

    async def generate_live_class_content(self, topic: str) -> Dict[str, Any]:
            """Generate Live Class content: MCQs, Polls, Flashcards"""
            try:
                print(f" [GEMINI_CODING] Generating Live Class content for topic: {topic}")
                
                if not self.available:
                    # Return basic fallback structure
                    return {
                        "summary": "Session summary unavailable.",
                        "quizzes": [],
                        "polls": [],
                        "flashcards": ["Fallback Flashcard 1", "Fallback Flashcard 2"]
                    }
                
                prompt = f"""
                Generate content for a Live Class session on the topic: "{topic}".
                
                Requirements:
                1. 5 MCQ Questions (Assessment)
                   - format: {{ "question": "...", "options": ["A","B","C","D"], "correct_option": 0 }}
                   - "correct_option" is index 0-3
                2. 3 Pulse Check Polls (Understanding check)
                   - format: {{ "text": "...", "type": "POLL", "options": ["Yes", "No", "Somewhat"] }}
                3. 5 Key Definition Flashcards
                   - format: simple string "Term: Definition"
                4. A short summary of the topic (approx 50 words)
                   - format: string "summary": "..."
                
                Return ONLY a valid JSON object with this EXACT structure:
                {{
                    "summary": "This session covers...",
                    "quizzes": [
                        {{
                            "title": "Quick Quiz",
                            "questions": [
                                {{ "text": "...", "type": "MCQ", "options": ["..."], "correct_option": 0 }}
                            ]
                        }}
                    ],
                    "polls": [
                        {{ "text": "...", "type": "POLL", "options": ["..."] }}
                    ],
                    "flashcards": [
                        "Term: Definition",
                        "..."
                    ]
                }}
                """
                
                try:
                    response = await self.client.aio.models.generate_content(
                        model=self.model_name,
                        contents=prompt,
                        config=types.GenerateContentConfig(
                            temperature=0.7,
                            max_output_tokens=4096
                        )
                    )
                except Exception as e:
                    print(f"[ERROR] [GEMINI_CODING] Live Content generation failed: {e}")
                    return {"quizzes": [], "polls": [], "flashcards": []}
                
                if not response.text:
                     return {"quizzes": [], "polls": [], "flashcards": []}
                     
                json_text = self._clean_json_response(response.text.strip())
                
                try:
                    data = json.loads(json_text)
                    
                    # Basic validation/cleanup can go here if needed
                    # Ensure "quizzes" structure matches what frontend expects for Assessment
                    # The prompt asks for "questions" list inside "quizzes"; we need to wrap it if needed or 
                    # in the prompt we asked for "quizzes" as a list of Assessments. 
                    # Actually, in LiveContent model: quizzes: List[Assessment]
                    # Assessment has: title, questions: List[Question]
                    
                    return data
                except json.JSONDecodeError:
                    print(f"[ERROR] [GEMINI_CODING] JSON parse failed for Live Content")
                    return {"quizzes": [], "polls": [], "flashcards": []}

            except Exception as e:
                print(f" [GEMINI_CODING] Error generating live content: {e}")
                return {"quizzes": [], "polls": [], "flashcards": []}

    async def generate_content_from_file(self, file_content: bytes, mime_type: str, topic: str = "General") -> Dict[str, Any]:
            """Generate Live Class content from an uploaded file (PDF/PPT/Text/Image)"""
            try:
                print(f" [GEMINI_CODING] Generating content from file ({mime_type}) for topic: {topic}")
                
                if not self.available:
                    print("[WARNING] [GEMINI_CODING] Service not available (missing API key?)")
                    return {
                        "summary": f"Content generation unavailable for {topic}.",
                        "quizzes": [],
                        "polls": [],
                        "flashcards": ["Feature Unavailable"]
                    }
                
                print(f"[DEBUG] [GEMINI_CODING] Input File Size: {len(file_content)} bytes")
                
                # Extract Text from File
                extracted_text = ""
                is_image = False
                
                try:
                    if mime_type == "application/pdf":
                        import io
                        import PyPDF2
                        pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
                        for page in pdf_reader.pages:
                            text = page.extract_text()
                            if text:
                                extracted_text += text + "\n"
                        # If empty, might be image-based PDF
                        if not extracted_text.strip():
                             extracted_text = "[PDF contains images or no selectable text]"

                    elif mime_type in ["application/vnd.openxmlformats-officedocument.presentationml.presentation", "application/vnd.ms-powerpoint"]:
                        import io
                        from pptx import Presentation
                        prs = Presentation(io.BytesIO(file_content))
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    extracted_text += shape.text + "\n"
                                    
                    elif mime_type.startswith("image/"):
                        is_image = True
                        extracted_text = "[Image File Provided]"
                        
                    else:
                        # Try generic text decode
                        extracted_text = file_content.decode("utf-8", errors="ignore")
                        
                except Exception as e:
                    print(f"[ERROR] Text extraction failed: {e}")
                    return {"quizzes": [], "polls": [], "flashcards": [], "fillups": [], "summary": f"Error parsing file: {str(e)}"}

                # Prepare Prompt
                prompt_instruction = f"""
                You are an expert educational content creator. 
                Analyze the following course material content covering the topic "{topic}":
                
                --- BEGIN CONTENT ---
                {extracted_text[:50000] if not is_image else "[Image Attached]"}
                --- END CONTENT ---
                
                Based on this content, generate the following interactive elements for a live class:
                
                1. **Today's Topic**: Extract the main topic name (e.g., "Photosynthesis", "Linear Algebra").
                2. **5 MCQ Questions** (for assessment)
                   - Must be directly answered by the file content.
                   - Format: {{ "question": "...", "options": ["A","B","C","D"], "correct_option": 0, "explanation": "..." }}
                3. **3 Pulse Check Polls** (to check understanding)
                   - Format: {{ "text": "...", "type": "POLL", "options": ["Yes", "No", "Somewhat"] }}
                4. **5 Flashcards** (Key terms/concepts)
                   - Format: "Term: Definition"
                5. **5 Fill-in-the-blank Questions** (for recall)
                   - Format: {{ "text": "The _______ is the powerhouse of the cell.", "answer": "mitochondria" }}
                6. **Brief Summary** (100 words, markdown supported)
                   - Provide a clear, structured summary of the file content.
                7. **1 Coding Problem** (Optional, only if technical)
                   - Format: {{ "title": "...", "description": "..." }} or null.
                
                Return ONLY a valid JSON object with this EXACT structure:
                {{
                    "summary": "### Key Concepts\\n...",
                    "quizzes": [
                        {{
                            "title": "Topic Quiz",
                            "questions": [
                                {{ "text": "...", "type": "MCQ", "options": ["..."], "correct_option": 0, "explanation": "..." }}
                            ]
                        }}
                    ],
                    "polls": [ {{ "text": "...", "type": "POLL", "options": ["..."] }} ],
                    "flashcards": [ "Term: Definition" ],
                    "fillups": [ {{ "text": "...", "answer": "..." }} ],
                    "coding_problem": {{ ... }}
                }}
                """
                
                contents = [prompt_instruction]
                
                if is_image:
                     contents.append({
                        "mime_type": mime_type,
                        "data": file_content
                    })
                
                import asyncio
                try:
                    # Use new SDK client.aio.models.generate_content for async
                    response = await self.client.aio.models.generate_content(
                        model=self.model_name,
                        contents=contents,
                        config=types.GenerateContentConfig(
                            temperature=0.7,
                            max_output_tokens=4096
                        )
                    )
                except Exception as e:
                    print(f"[ERROR] [GEMINI_CODING] File processing failed: {e}")
                    return {"quizzes": [], "polls": [], "flashcards": [], "fillups": [], "summary": "Error: Analysis failed."}
                    
                if not response.text:
                    return {"quizzes": [], "polls": [], "flashcards": [], "fillups": [], "summary": "Empty response from AI."}
                    
                json_text = self._clean_json_response(response.text.strip())
                try:
                    data = json.loads(json_text)
                    return data
                except json.JSONDecodeError:
                    print(f"[ERROR] JSON parse failed: {json_text[:200]}")
                    return {"quizzes": [], "polls": [], "flashcards": [], "fillups": [], "summary": "Error parsing AI response."}

            except Exception as e:
                print(f" [GEMINI_CODING] Error: {e}")
                return {"quizzes": [], "polls": [], "flashcards": [], "fillups": [], "summary": f"System Error: {str(e)}"}

    def _get_fallback_coding_mcq(self, topic: str, difficulty: str) -> Dict[str, Any]:
            """Get a fallback coding MCQ when AI is not available"""
            import random
            
            # Coding MCQ questions database
            coding_mcqs = {
                "Data Structures": {
                    "easy": [
                        {
                            "question": "What is the time complexity of accessing an element in an array?",
                            "options": ["O(1)", "O(n)", "O(log n)", "O(n)"],
                            "correct_answer": 0,
                            "explanation": "Array access is O(1) because we can directly access any element using its index."
                        },
                        {
                            "question": "Which data structure follows LIFO (Last In, First Out) principle?",
                            "options": ["Queue", "Stack", "Array", "Linked List"],
                            "correct_answer": 1,
                            "explanation": "Stack follows LIFO principle where the last element added is the first one to be removed."
                        }
                    ],
                    "medium": [
                        {
                            "question": "What is the time complexity of inserting an element in a balanced binary search tree?",
                            "options": ["O(1)", "O(log n)", "O(n)", "O(n log n)"],
                            "correct_answer": 1,
                            "explanation": "In a balanced BST, insertion requires traversing the height of the tree, which is O(log n)."
                        }
                    ],
                    "hard": [
                        {
                            "question": "What is the space complexity of merge sort?",
                            "options": ["O(1)", "O(log n)", "O(n)", "O(n log n)"],
                            "correct_answer": 2,
                            "explanation": "Merge sort requires O(n) extra space for the temporary arrays used during merging."
                        }
                    ]
                },
                "Algorithms": {
                    "easy": [
                        {
                            "question": "Which sorting algorithm has the best average-case time complexity?",
                            "options": ["Bubble Sort", "Selection Sort", "Quick Sort", "Insertion Sort"],
                            "correct_answer": 2,
                            "explanation": "Quick Sort has O(n log n) average-case time complexity, which is better than the O(n) of the others."
                        }
                    ],
                    "medium": [
                        {
                            "question": "What is the time complexity of Dijkstra's algorithm?",
                            "options": ["O(V)", "O(V + E)", "O(V log V + E)", "O(V)"],
                            "correct_answer": 2,
                            "explanation": "Dijkstra's algorithm with a binary heap has O(V log V + E) time complexity."
                        }
                    ],
                    "hard": [
                        {
                            "question": "What is the space complexity of recursive Fibonacci without memoization?",
                            "options": ["O(1)", "O(n)", "O(2^n)", "O(log n)"],
                            "correct_answer": 2,
                            "explanation": "Recursive Fibonacci without memoization has exponential time and space complexity O(2^n)."
                        }
                    ]
                },
                "Python Programming": {
                    "easy": [
                        {
                            "question": "What is the output of: print(3 * 'abc')",
                            "options": ["abcabcabc", "3abc", "abc3", "Error"],
                            "correct_answer": 0,
                            "explanation": "Multiplying a string by an integer repeats the string that many times."
                        }
                    ],
                    "medium": [
                        {
                            "question": "What does the 'with' statement in Python provide?",
                            "options": ["Loop control", "Exception handling", "Resource management", "Function definition"],
                            "correct_answer": 2,
                            "explanation": "The 'with' statement provides automatic resource management and cleanup."
                        }
                    ],
                    "hard": [
                        {
                            "question": "What is the difference between 'is' and '==' in Python?",
                            "options": ["No difference", "'is' compares values, '==' compares identity", "'is' compares identity, '==' compares values", "Both compare identity"],
                            "correct_answer": 2,
                            "explanation": "'is' checks if two variables refer to the same object (identity), while '==' checks if they have the same value."
                        }
                    ]
                }
            }
            
            # Get questions for the topic and difficulty
            if topic in coding_mcqs and difficulty in coding_mcqs[topic]:
                questions = coding_mcqs[topic][difficulty]
                if isinstance(questions, list):
                    selected_question = random.choice(questions)
                else:
                    selected_question = questions
            else:
                # Default fallback question
                selected_question = {
                    "question": "What is the time complexity of linear search?",
                    "options": ["O(1)", "O(log n)", "O(n)", "O(n)"],
                    "correct_answer": 2,
                    "explanation": "Linear search checks each element one by one, so it has O(n) time complexity in the worst case."
                }
            
            # Add topic and difficulty
            selected_question['topic'] = topic
            selected_question['difficulty'] = difficulty
            
            return selected_question

